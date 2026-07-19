#!/usr/bin/env python3
"""Dependency-free regression tests for the bundled scripts.

Run from the repo root:
    python3 -m unittest discover -s tests -v

Uses only the standard library (unittest). Pure-function behaviour is tested by
importing the scripts directly; CLI/exit-code behaviour (validate, importers) is
tested via subprocess against tiny temp fixtures. No Graphviz/draw.io needed —
auto-layout is exercised through to_drawio() with synthetic node positions.
"""
import base64
import gzip
import importlib.util
import json
import os
import subprocess
import sys
import tempfile
import unittest
import urllib.parse
import zlib

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SCRIPTS = os.path.join(ROOT, "skills", "drawio-skill", "scripts")
DATA = os.path.join(ROOT, "skills", "drawio-skill", "data")


def load(name):
    """Import a bundled script module by file path."""
    path = os.path.join(SCRIPTS, name + ".py")
    spec = importlib.util.spec_from_file_location(name, path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def run(script, *args, **kw):
    """Run a script as a subprocess; return CompletedProcess."""
    return subprocess.run([sys.executable, os.path.join(SCRIPTS, script), *args],
                          capture_output=True, text=True, **kw)


class TestShapeSearch(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.m = load("shapesearch")
        with gzip.open(os.path.join(DATA, "shape-index.json.gz"), "rt", encoding="utf-8") as f:
            cls.shapes = json.load(f)
        cls.tm = cls.m.build_tag_map(cls.shapes)

    def search(self, q, n=5):
        return self.m.search(self.shapes, self.tm, q, n)

    def test_soundex(self):
        self.assertEqual(self.m.soundex("Robert"), "R163")
        self.assertEqual(self.m.soundex("Jackson"), "J250")

    def test_index_loaded(self):
        self.assertGreater(len(self.shapes), 10000)
        self.assertTrue(self.tm)

    def test_known_shapes(self):
        self.assertIn("Lambda", self.search("aws lambda")[0]["title"])
        self.assertIn("Actor", self.search("uml actor")[0]["title"])

    def test_title_exact_ranking(self):
        # The v1.11.1 fix: the shape literally titled "DynamoDB" ranks first,
        # above tag-only neighbours like "Attribute".
        self.assertEqual(self.search("aws dynamodb")[0]["title"], "DynamoDB")

    def test_no_match_is_empty(self):
        self.assertEqual(self.search("zzzznotashape"), [])


class TestAiIcons(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.m = load("aiicons")
        with open(os.path.join(DATA, "lobe-icons.json"), encoding="utf-8") as fh:
            manifest = json.load(fh)
        cls.fam = cls.m.families(manifest["icons"])

    def test_families_group_variants(self):
        self.assertIn("claude", self.fam)
        self.assertIn("claude-color", self.fam["claude"])
        self.assertIn("googlecloud-brand", self.fam["googlecloud"])
        self.assertIn("alibabacloud-text-cn", self.fam["alibabacloud"])
        self.assertNotIn("googlecloud-brand", self.fam)
        self.assertNotIn("alibabacloud-text-cn", self.fam)

    def test_search_matches_brand(self):
        self.assertEqual(self.m.search(self.fam, "claude", 1)[0], "claude")
        # token matching: a brand word inside a phrase still matches
        self.assertEqual(self.m.search(self.fam, "use the openai logo", 1)[0], "openai")
        self.assertEqual(self.m.search(self.fam, "Atlas Cloud", 1)[0], "atlascloud")

    def test_search_does_not_return_variant_brands(self):
        matches = self.m.search(self.fam, "google cloud", 8)
        self.assertIn("googlecloud", matches)
        self.assertNotIn("googlecloud-brand", matches)
        matches = self.m.search(self.fam, "alibaba cloud", 8)
        self.assertIn("alibabacloud", matches)
        self.assertNotIn("alibabacloud-text-cn", matches)

    def test_variant_preference(self):
        # claude has a colour variant; openai is mono-only -> falls back.
        self.assertEqual(self.m.pick_variant("claude", self.fam["claude"], "color"), "claude-color")
        self.assertEqual(self.m.pick_variant("openai", self.fam["openai"], "color"), "openai")
        self.assertEqual(self.m.pick_variant("googlecloud", self.fam["googlecloud"], "text"),
                         "googlecloud-brand")

    def test_unknown_brand(self):
        self.assertEqual(self.m.search(self.fam, "definitelynotabrand", 3), [])


class TestEncodeUrl(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.m = load("encode_drawio_url")

    @staticmethod
    def decode(b64):
        # Reverse the encoder: inflate, then decodeURIComponent.
        inflated = zlib.decompress(base64.b64decode(b64), -zlib.MAX_WBITS).decode("utf-8")
        return urllib.parse.unquote(inflated)

    def test_roundtrip_cjk_and_percent(self):
        # The v1.11.0 fix: encodeURIComponent before deflate keeps % and CJK safe.
        xml = '<mxGraphModel><root><mxCell value="登录 100% &amp; ok"/></root></mxGraphModel>'
        b64 = self.m._deflate_b64(xml)
        self.assertEqual(self.decode(b64), xml)

    def test_viewer_and_editor_urls(self):
        xml = "<mxGraphModel><root/></mxGraphModel>"
        self.assertIn("viewer.diagrams.net", self.m.encode(xml))
        self.assertIn("#R", self.m.encode(xml))
        self.assertIn("app.diagrams.net", self.m.edit_url(xml))
        self.assertIn("#create=", self.m.edit_url(xml))


class TestAutolayoutColor(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.m = load("autolayout")
        cls.pal = cls.m.load_palette()

    def test_palette_from_design_file(self):
        # Sourced from styles/built-in/default.json, not the inline fallback.
        self.assertNotEqual(self.pal, self.m._FALLBACK_PALETTE)
        self.assertEqual(self.pal[0], ("#dae8fc", "#6c8ebf"))   # primary

    def _render(self, color):
        graph = {"direction": "TB", "nodes": [
            {"id": "a", "label": "A", "group": "mod1"},
            {"id": "b", "label": "B", "group": "mod2"},
            {"id": "c", "label": "C", "style": "fillColor=#ff0000;"},
            {"id": "d", "label": "D"},
        ], "edges": []}
        pos = {"a": (1, 4), "b": (3, 4), "c": (5, 4), "d": (7, 4)}
        return self.m.to_drawio(graph, 5, pos, {}, color=color)

    def test_grouped_nodes_coloured(self):
        out = self._render(color=True)
        self.assertIn("fillColor=#d5e8d4", out)   # b tinted by mod2 (success/green)
        self.assertIn("strokeColor=#82b366", out)  # mod2 container border colour
        self.assertIn("fillColor=#ff0000", out)    # explicit style preserved

    def test_mono_disables_colour(self):
        out = self._render(color=False)
        self.assertIn("strokeColor=#999999", out)  # grey container
        self.assertNotIn("fillColor=#d5e8d4", out)  # no group tint
        self.assertIn("fillColor=#ff0000", out)     # explicit style still wins

    def test_graph_level_spacing(self):
        # Importers emitting icon nodes ask for extra separation (labels render
        # below the shape); the keys pass straight through to dot.
        dot = self.m.build_dot({"nodes": [{"id": "a"}], "edges": [],
                                "ranksep": 0.7, "nodesep": 0.6})
        self.assertIn("ranksep=0.70;", dot)
        self.assertIn("nodesep=0.60;", dot)

    def test_node_link_wraps_userobject(self):
        out = self.m.to_drawio(
            {"direction": "TB",
             "nodes": [{"id": "n", "label": "L", "link": "data:page/id,p2"}],
             "edges": []}, 5, {"n": (1, 1)}, {}, color=True)
        self.assertIn('<UserObject label="L" link="data:page/id,p2" id="n">', out)
        self.assertNotIn('value="L"', out)          # label moved to the wrapper

    def test_edge_style_passthrough(self):
        out = self.m.to_drawio(
            {"direction": "TB",
             "nodes": [{"id": "a"}, {"id": "b"}],
             "edges": [{"source": "a", "target": "b", "style": "endArrow=ERone;"}]},
            5, {"a": (1, 1), "b": (3, 1)}, {}, color=True)
        self.assertIn("endArrow=ERone;", out)

    def test_label_newline_becomes_entity(self):
        out = self.m.to_drawio(
            {"direction": "TB", "nodes": [{"id": "n", "label": "one\ntwo"}],
             "edges": []}, 5, {"n": (1, 1)}, {}, color=True)
        self.assertIn("one&#xa;two", out)

    def test_escapes_special_chars(self):
        # Ids/edges with a quote or backslash must be escaped in the DOT input
        # (else they corrupt the Graphviz source); a style with a quote must be
        # XML-escaped in the emitted mxCell.
        dot = self.m.build_dot({
            "nodes": [{"id": 'a"x'}, {"id": "b\\y"}],
            "edges": [{"source": 'a"x', "target": "b\\y"}],
        })
        self.assertIn(r'"a\"x"', dot)
        self.assertIn(r'"b\\y"', dot)
        out = self.m.to_drawio(
            {"direction": "TB",
             "nodes": [{"id": "n", "label": "L", "style": 'fillColor="#fff";'}],
             "edges": []},
            5, {"n": (1, 1)}, {}, color=True)
        self.assertIn("&quot;", out)


class TestBuiltinPresets(unittest.TestCase):
    STYLES = os.path.join(ROOT, "skills", "drawio-skill", "styles", "built-in")

    def test_presets_conform(self):
        import glob
        import re
        hexpat = re.compile(r"^#[0-9A-Fa-f]{6}$")
        files = sorted(glob.glob(os.path.join(self.STYLES, "*.json")))
        names = {os.path.basename(f)[:-5] for f in files}
        self.assertLessEqual({"default", "corporate", "handdrawn",
                              "colorblind-safe", "dark"}, names)
        for f in files:
            with open(f, encoding="utf-8") as fh:
                d = json.load(fh)
            self.assertEqual(d["name"], os.path.basename(f)[:-5])
            self.assertFalse(d.get("default"), f)   # built-ins never default
            for slot in ("primary", "success", "warning", "accent",
                         "danger", "neutral", "secondary"):
                pair = d["palette"][slot]
                if pair is not None:
                    self.assertRegex(pair["fillColor"], hexpat)
                    self.assertRegex(pair["strokeColor"], hexpat)
            for key in ("background", "fontColor", "edgeColor"):
                if key in d.get("extras", {}):
                    self.assertRegex(d["extras"][key], hexpat)


class TestValidateCli(unittest.TestCase):
    GOOD = ('<mxfile><diagram name="P1"><mxGraphModel><root>'
            '<mxCell id="0"/><mxCell id="1" parent="0"/>'
            '<mxCell id="2" value="A" vertex="1" parent="1">'
            '<mxGeometry x="0" y="0" width="80" height="40" as="geometry"/></mxCell>'
            '<mxCell id="3" value="B" vertex="1" parent="1">'
            '<mxGeometry x="200" y="0" width="80" height="40" as="geometry"/></mxCell>'
            '<mxCell id="4" edge="1" parent="1" source="2" target="3">'
            '<mxGeometry relative="1" as="geometry"/></mxCell>'
            '</root></mxGraphModel></diagram></mxfile>')
    BAD = ('<mxfile><diagram name="P1"><mxGraphModel><root>'
           '<mxCell id="0"/><mxCell id="1" parent="0"/>'
           '<mxCell id="2" value="A" vertex="1" parent="1">'
           '<mxGeometry x="0" y="0" width="80" height="40" as="geometry"/></mxCell>'
           '<mxCell id="4" edge="1" parent="1" source="2" target="nope">'
           '<mxGeometry relative="1" as="geometry"/></mxCell>'
           '</root></mxGraphModel></diagram></mxfile>')
    # Edge label: a relative-positioned vertex that legitimately omits width/height.
    EDGE_LABEL = ('<mxfile><diagram name="P1"><mxGraphModel><root>'
                  '<mxCell id="0"/><mxCell id="1" parent="0"/>'
                  '<mxCell id="2" value="A" vertex="1" parent="1">'
                  '<mxGeometry x="0" y="0" width="80" height="40" as="geometry"/></mxCell>'
                  '<mxCell id="3" value="B" vertex="1" parent="1">'
                  '<mxGeometry x="200" y="0" width="80" height="40" as="geometry"/></mxCell>'
                  '<mxCell id="4" edge="1" parent="1" source="2" target="3">'
                  '<mxGeometry relative="1" as="geometry"/></mxCell>'
                  '<mxCell id="5" value="lbl" style="edgeLabel;html=1;" vertex="1" '
                  'connectable="0" parent="1">'
                  '<mxGeometry x="0.5" y="0" relative="1" as="geometry">'
                  '<mxPoint as="offset"/></mxGeometry></mxCell>'
                  '</root></mxGraphModel></diagram></mxfile>')

    def _check(self, xml):
        with tempfile.NamedTemporaryFile("w", suffix=".drawio", delete=False) as f:
            f.write(xml)
            path = f.name
        try:
            return run("validate.py", path)
        finally:
            os.unlink(path)

    def test_good_passes(self):
        self.assertEqual(self._check(self.GOOD).returncode, 0)

    def test_dangling_edge_fails(self):
        r = self._check(self.BAD)
        self.assertEqual(r.returncode, 1)
        self.assertIn("error", r.stdout)

    # A vertex wrapped in UserObject (link) — the edge references the wrapper id.
    LINKED = ('<mxfile><diagram name="P1"><mxGraphModel><root>'
              '<mxCell id="0"/><mxCell id="1" parent="0"/>'
              '<UserObject label="A" link="data:page/id,p2" id="2">'
              '<mxCell vertex="1" parent="1">'
              '<mxGeometry x="0" y="0" width="80" height="40" as="geometry"/></mxCell>'
              '</UserObject>'
              '<mxCell id="3" value="B" vertex="1" parent="1">'
              '<mxGeometry x="200" y="0" width="80" height="40" as="geometry"/></mxCell>'
              '<mxCell id="4" edge="1" parent="1" source="2" target="3">'
              '<mxGeometry relative="1" as="geometry"/></mxCell>'
              '</root></mxGraphModel></diagram></mxfile>')

    def test_userobject_wrapper_resolves(self):
        # The wrapper id must satisfy edge source refs (no dangling-edge error).
        r = self._check(self.LINKED)
        self.assertEqual(r.returncode, 0, r.stdout)
        self.assertIn("0 error(s)", r.stdout)

    def test_edge_label_passes(self):
        # edgeLabel vertices have no width/height; they must not be flagged
        # as missing/invalid geometry (issue #35).
        r = self._check(self.EDGE_LABEL)
        self.assertEqual(r.returncode, 0)
        self.assertNotIn("error:", r.stdout)


class TestValidateGeometry(unittest.TestCase):
    """Edge routing geometry checks (waypointed edges only)."""

    @classmethod
    def setUpClass(cls):
        cls.m = load("validate")

    @staticmethod
    def _ids(xml):
        import xml.etree.ElementTree as ET
        root = ET.fromstring(xml)
        return {c.get("id"): c for c in root.iter("mxCell")}

    # --- pure-function unit tests ---
    def test_segments_cross(self):
        c = self.m.segments_cross
        self.assertTrue(c((0, 0), (10, 10), (0, 10), (10, 0)))     # an X
        self.assertFalse(c((0, 0), (10, 0), (0, 5), (10, 5)))      # parallel
        self.assertFalse(c((0, 0), (10, 0), (10, 0), (10, 10)))    # touch at endpoint

    def test_route_hits_rect(self):
        h = self.m.route_hits_rect
        self.assertTrue(h([(0, 50), (100, 50)], (40, 40, 20, 20)))   # cuts through
        self.assertFalse(h([(0, 0), (100, 0)], (40, 40, 20, 20)))    # clear of box

    def test_abs_rect_resolves_container_offset(self):
        ids = self._ids(
            '<root>'
            '<mxCell id="c" vertex="1" parent="1">'
            '<mxGeometry x="100" y="100" width="200" height="200" as="geometry"/></mxCell>'
            '<mxCell id="ch" vertex="1" parent="c">'
            '<mxGeometry x="10" y="10" width="40" height="40" as="geometry"/></mxCell>'
            '</root>')
        self.assertEqual(self.m.abs_rect(ids["ch"], ids), (110, 110, 40, 40))

    def test_endpoint_honours_exit_point(self):
        ids = self._ids(
            '<root>'
            '<mxCell id="s" vertex="1" parent="1">'
            '<mxGeometry x="0" y="0" width="80" height="40" as="geometry"/></mxCell>'
            '<mxCell id="e" edge="1" parent="1" source="s" target="s" '
            'style="exitX=1;exitY=0.5;"><mxGeometry relative="1" as="geometry"/></mxCell>'
            '</root>')
        self.assertEqual(self.m.endpoint(ids["e"], "source", ids), (80, 20))

    # --- CLI integration tests ---
    def _check(self, xml):
        with tempfile.NamedTemporaryFile("w", suffix=".drawio", delete=False) as f:
            f.write(xml)
            path = f.name
        try:
            return run("validate.py", path)
        finally:
            os.unlink(path)

    # Two waypointed edges forming an X that crosses at (300,300).
    CROSS = ('<mxfile><diagram name="P1"><mxGraphModel><root>'
             '<mxCell id="0"/><mxCell id="1" parent="0"/>'
             '<mxCell id="n1" vertex="1" parent="1"><mxGeometry x="80" y="80" width="40" height="40" as="geometry"/></mxCell>'
             '<mxCell id="n2" vertex="1" parent="1"><mxGeometry x="80" y="480" width="40" height="40" as="geometry"/></mxCell>'
             '<mxCell id="n3" vertex="1" parent="1"><mxGeometry x="480" y="80" width="40" height="40" as="geometry"/></mxCell>'
             '<mxCell id="n4" vertex="1" parent="1"><mxGeometry x="480" y="480" width="40" height="40" as="geometry"/></mxCell>'
             '<mxCell id="e1" edge="1" parent="1" source="n1" target="n4"><mxGeometry relative="1" as="geometry">'
             '<Array as="points"><mxPoint x="200" y="200"/></Array></mxGeometry></mxCell>'
             '<mxCell id="e2" edge="1" parent="1" source="n2" target="n3"><mxGeometry relative="1" as="geometry">'
             '<Array as="points"><mxPoint x="200" y="400"/></Array></mxGeometry></mxCell>'
             '</root></mxGraphModel></diagram></mxfile>')

    # A waypointed edge whose route passes straight through obstacle vertex 'v'.
    THROUGH = ('<mxfile><diagram name="P1"><mxGraphModel><root>'
               '<mxCell id="0"/><mxCell id="1" parent="0"/>'
               '<mxCell id="n1" vertex="1" parent="1"><mxGeometry x="80" y="80" width="40" height="40" as="geometry"/></mxCell>'
               '<mxCell id="n2" vertex="1" parent="1"><mxGeometry x="480" y="80" width="40" height="40" as="geometry"/></mxCell>'
               '<mxCell id="v" vertex="1" parent="1"><mxGeometry x="280" y="80" width="40" height="40" as="geometry"/></mxCell>'
               '<mxCell id="e" edge="1" parent="1" source="n1" target="n2"><mxGeometry relative="1" as="geometry">'
               '<Array as="points"><mxPoint x="300" y="100"/></Array></mxGeometry></mxCell>'
               '</root></mxGraphModel></diagram></mxfile>')

    # Same geometry as THROUGH but the edge has NO waypoints — auto-routed, so
    # its path is unknown and must NOT be geometry-checked (no false positive).
    AUTOROUTE = ('<mxfile><diagram name="P1"><mxGraphModel><root>'
                 '<mxCell id="0"/><mxCell id="1" parent="0"/>'
                 '<mxCell id="n1" vertex="1" parent="1"><mxGeometry x="80" y="80" width="40" height="40" as="geometry"/></mxCell>'
                 '<mxCell id="n2" vertex="1" parent="1"><mxGeometry x="480" y="80" width="40" height="40" as="geometry"/></mxCell>'
                 '<mxCell id="v" vertex="1" parent="1"><mxGeometry x="280" y="80" width="40" height="40" as="geometry"/></mxCell>'
                 '<mxCell id="e" edge="1" parent="1" source="n1" target="n2">'
                 '<mxGeometry relative="1" as="geometry"/></mxCell>'
                 '</root></mxGraphModel></diagram></mxfile>')

    def test_edge_crossing_warns(self):
        r = self._check(self.CROSS)
        self.assertIn("edges 'e1' and 'e2' cross", r.stdout)
        self.assertEqual(r.returncode, 0)            # warning, not error

    def test_edge_crossing_strict_fails(self):
        with tempfile.NamedTemporaryFile("w", suffix=".drawio", delete=False) as f:
            f.write(self.CROSS)
            path = f.name
        try:
            self.assertEqual(run("validate.py", path, "--strict").returncode, 1)
        finally:
            os.unlink(path)

    def test_edge_through_vertex_warns(self):
        r = self._check(self.THROUGH)
        self.assertIn("edge 'e' routes through vertex 'v'", r.stdout)

    def test_autorouted_edge_not_checked(self):
        # No waypoints -> path unknown -> no geometry warning (FP-free).
        r = self._check(self.AUTOROUTE)
        self.assertNotIn("routes through", r.stdout)
        self.assertEqual(r.returncode, 0)


class TestImportersCli(unittest.TestCase):
    @staticmethod
    def _write(path, text):
        with open(path, "w", encoding="utf-8") as f:
            f.write(text)

    def test_pyimports_edge(self):
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "a.py"), "import b\n")
            self._write(os.path.join(d, "b.py"), "x = 1\n")
            r = run("pyimports.py", d)
            graph = json.loads(r.stdout)
            ids = {n["id"] for n in graph["nodes"]}
            self.assertEqual(ids, {"a", "b"})
            self.assertIn({"source": "a", "target": "b"}, graph["edges"])

    def test_pyclasses_inheritance(self):
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "m.py"), "class A: pass\nclass B(A): pass\n")
            r = run("pyclasses.py", d)
            graph = json.loads(r.stdout)
            self.assertEqual(len(graph["nodes"]), 2)
            self.assertEqual(graph["edges"], [{"source": "m.B", "target": "m.A"}])
            # No hard-coded colour anymore (coloured by group in autolayout).
            self.assertNotIn("style", graph["nodes"][0])

    def test_jsimports_edge(self):
        # Node ids are the module basename without extension: a, b.
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "a.js"), "import './b.js';\n")
            self._write(os.path.join(d, "b.js"), "export const x = 1;\n")
            r = run("jsimports.py", d)
            graph = json.loads(r.stdout)
            ids = {n["id"] for n in graph["nodes"]}
            self.assertEqual(len(graph["nodes"]), 2)
            self.assertEqual(ids, {"a", "b"})
            self.assertEqual(len(graph["edges"]), 1)
            self.assertEqual(graph["edges"][0], {"source": "a", "target": "b"})

    def test_goimports_edge(self):
        # Node ids are full import paths: example.com/m/<pkg>.
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "go.mod"), "module example.com/m\n\ngo 1.21\n")
            os.mkdir(os.path.join(d, "a"))
            os.mkdir(os.path.join(d, "b"))
            self._write(os.path.join(d, "a", "a.go"),
                        'package a\n\nimport _ "example.com/m/b"\n')
            self._write(os.path.join(d, "b", "b.go"), "package b\n")
            r = run("goimports.py", d)
            graph = json.loads(r.stdout)
            ids = {n["id"] for n in graph["nodes"]}
            self.assertEqual(len(graph["nodes"]), 2)
            self.assertEqual(ids, {"example.com/m/a", "example.com/m/b"})
            self.assertEqual(len(graph["edges"]), 1)
            self.assertEqual(graph["edges"][0],
                             {"source": "example.com/m/a", "target": "example.com/m/b"})

    def test_tfimports_edge_and_icons(self):
        # Node ids are type.name; the lambda resolves to an official aws4 icon
        # and the referenced role becomes an edge.
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "main.tf"),
                        'resource "aws_lambda_function" "fn" {\n'
                        '  role = aws_iam_role.exec.arn  # comment\n'
                        '}\n'
                        'resource "aws_iam_role" "exec" {\n'
                        '  name = "exec"\n'
                        '}\n')
            r = run("tfimports.py", d, "--no-reduce")
            graph = json.loads(r.stdout)
            ids = {n["id"] for n in graph["nodes"]}
            self.assertEqual(ids, {"aws_lambda_function.fn", "aws_iam_role.exec"})
            self.assertEqual(graph["edges"], [{"source": "aws_lambda_function.fn",
                                               "target": "aws_iam_role.exec"}])
            fn = next(n for n in graph["nodes"] if n["id"] == "aws_lambda_function.fn")
            self.assertIn("mxgraph.aws4", fn["style"])
            self.assertEqual(graph["ranksep"], 0.7)   # icon labels need spacing

    def test_tfimports_no_icons(self):
        # Without icons the type stays visible on the box (two-line label).
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "main.tf"),
                        'resource "aws_s3_bucket" "logs" {}\n')
            graph = json.loads(run("tfimports.py", d, "--no-icons").stdout)
            self.assertEqual(graph["nodes"][0]["label"], "logs\naws_s3_bucket")
            self.assertNotIn("style", graph["nodes"][0])

    K8S_LIST = {"kind": "List", "items": [
        {"apiVersion": "v1", "kind": "Service",
         "metadata": {"name": "web", "namespace": "shop"},
         "spec": {"selector": {"app": "web"}}},
        {"apiVersion": "apps/v1", "kind": "Deployment",
         "metadata": {"name": "web", "namespace": "shop"},
         "spec": {"template": {"metadata": {"labels": {"app": "web", "tier": "fe"}},
                  "spec": {"containers": [
                      {"name": "web", "image": "nginx",
                       "envFrom": [{"configMapRef": {"name": "cfg"}}]}]}}}},
        {"apiVersion": "v1", "kind": "ConfigMap",
         "metadata": {"name": "cfg", "namespace": "shop"}},
    ]}

    def test_k8simports_json(self):
        # JSON input needs no PyYAML; selector match and configMap ref -> edges.
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "all.json"), json.dumps(self.K8S_LIST))
            graph = json.loads(run("k8simports.py", d).stdout)
            ids = {n["id"] for n in graph["nodes"]}
            self.assertEqual(ids, {"shop/Service/web", "shop/Deployment/web",
                                   "shop/ConfigMap/cfg"})
            edges = {(e["source"], e["target"]) for e in graph["edges"]}
            self.assertEqual(edges, {("shop/Service/web", "shop/Deployment/web"),
                                     ("shop/Deployment/web", "shop/ConfigMap/cfg")})
            svc = next(n for n in graph["nodes"] if n["id"] == "shop/Service/web")
            self.assertIn("mxgraph.kubernetes", svc["style"])

    def test_k8simports_group_by_namespace(self):
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "all.json"), json.dumps(self.K8S_LIST))
            graph = json.loads(run("k8simports.py", d, "--group").stdout)
            self.assertTrue(all(n["group"] == "shop" for n in graph["nodes"]))

    def test_composeimports_edges(self):
        try:
            import yaml  # noqa: F401
        except ImportError:
            self.skipTest("PyYAML not installed")
        compose = ("services:\n"
                   "  api:\n    image: api:1\n    depends_on: [db]\n"
                   "    volumes: ['uploads:/data', './src:/app']\n"
                   "  db:\n    image: postgres:16\n"
                   "volumes:\n  uploads:\n")
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "docker-compose.yml"), compose)
            graph = json.loads(run("composeimports.py", d).stdout)
            ids = {n["id"] for n in graph["nodes"]}
            self.assertEqual(ids, {"api", "db", "vol:uploads"})
            edges = {(e["source"], e["target"]) for e in graph["edges"]}
            # bind mount ./src is not a named volume -> no node, no edge
            self.assertEqual(edges, {("api", "db"), ("api", "vol:uploads")})
            api = next(n for n in graph["nodes"] if n["id"] == "api")
            self.assertEqual(api["label"], "api\napi:1")

    DOCKER = [
        {"Id": "a1", "Name": "/shop-web-1",
         "Config": {"Image": "nginx:1.27", "Labels": {
             "com.docker.compose.project": "shop", "com.docker.compose.service": "web",
             "com.docker.compose.depends_on": "db:service_healthy:false"}},
         "NetworkSettings": {"Networks": {"shop_net": {}, "bridge": {}}}, "Mounts": []},
        {"Id": "a2", "Name": "/shop-db-1",
         "Config": {"Image": "postgres:16", "Labels": {
             "com.docker.compose.project": "shop", "com.docker.compose.service": "db"}},
         "NetworkSettings": {"Networks": {"shop_net": {}}},
         "Mounts": [{"Type": "volume", "Name": "pgdata", "Destination": "/data"},
                    {"Type": "bind", "Source": "/etc/hosts", "Destination": "/etc/hosts"}]},
    ]

    def test_dockerimports_live(self):
        # `docker inspect` array on stdin (-): depends_on resolves service->container;
        # built-in `bridge` network and the bind mount are dropped as noise.
        graph = json.loads(run("dockerimports.py", "-",
                               input=json.dumps(self.DOCKER)).stdout)
        ids = {n["id"] for n in graph["nodes"]}
        self.assertEqual(ids, {"shop-web-1", "shop-db-1", "net:shop_net", "vol:pgdata"})
        edges = {(e["source"], e["target"]) for e in graph["edges"]}
        self.assertEqual(edges, {("shop-web-1", "shop-db-1"),
                                 ("shop-web-1", "net:shop_net"),
                                 ("shop-db-1", "net:shop_net"),
                                 ("shop-db-1", "vol:pgdata")})
        web = next(n for n in graph["nodes"] if n["id"] == "shop-web-1")
        self.assertEqual(web["label"], "shop-web-1\nnginx:1.27")

    def test_dockerimports_group_by_project(self):
        graph = json.loads(run("dockerimports.py", "-", "--group",
                               input=json.dumps(self.DOCKER)).stdout)
        containers = [n for n in graph["nodes"] if not n["id"].startswith(("net:", "vol:"))]
        self.assertTrue(all(n["group"] == "shop" for n in containers))

    TFSTATE = {"format_version": "1.0", "values": {"root_module": {
        "resources": [
            {"address": "aws_instance.web", "mode": "managed", "type": "aws_instance",
             "name": "web", "depends_on": ["aws_security_group.web"], "values": {}},
            {"address": "aws_security_group.web", "mode": "managed",
             "type": "aws_security_group", "name": "web", "values": {}},
            {"address": "data.aws_ami.x", "mode": "data", "type": "aws_ami",
             "name": "x", "values": {}}],
        "child_modules": [{"address": "module.vpc", "resources": [
            {"address": "module.vpc.aws_subnet.p[0]", "mode": "managed",
             "type": "aws_subnet", "name": "p", "index": 0, "values": {}},
            {"address": "module.vpc.aws_subnet.p[1]", "mode": "managed",
             "type": "aws_subnet", "name": "p", "index": 1, "values": {}}]}]}}}

    def test_tfstate_live(self):
        # `terraform show -json` on stdin: data sources excluded, depends_on edge,
        # count instances expanded, aws_instance -> official icon, module grouping.
        graph = json.loads(run("tfstate.py", "-", "--group", "--no-reduce",
                               input=json.dumps(self.TFSTATE)).stdout)
        ids = {n["id"] for n in graph["nodes"]}
        self.assertEqual(ids, {"aws_instance.web", "aws_security_group.web",
                               "module.vpc.aws_subnet.p[0]", "module.vpc.aws_subnet.p[1]"})
        self.assertEqual(graph["edges"],
                         [{"source": "aws_instance.web", "target": "aws_security_group.web"}])
        ec2 = next(n for n in graph["nodes"] if n["id"] == "aws_instance.web")
        self.assertIn("mxgraph.aws4", ec2["style"])
        sub = next(n for n in graph["nodes"] if n["id"].endswith("p[0]"))
        self.assertEqual(sub["label"], "p[0]\naws_subnet")   # no icon -> type on box
        self.assertEqual(sub["group"], "module.vpc")

    def test_k8simports_stdin(self):
        # `kubectl get ... -o json | k8simports.py -` reads the live snapshot.
        graph = json.loads(run("k8simports.py", "-",
                               input=json.dumps(self.K8S_LIST)).stdout)
        self.assertEqual({n["id"] for n in graph["nodes"]},
                         {"shop/Service/web", "shop/Deployment/web", "shop/ConfigMap/cfg"})

    @staticmethod
    def _drawio(nodes, edges):
        """Minimal .drawio XML from (id, label) nodes and (src, tgt) edges."""
        cells = ['<mxCell id="0"/>', '<mxCell id="1" parent="0"/>']
        for i, (cid, label) in enumerate(nodes):
            cells.append(f'<mxCell id="{cid}" value="{label}" vertex="1" parent="1" '
                         f'style="rounded=1;"><mxGeometry x="{i*200}" y="0" '
                         'width="120" height="60" as="geometry"/></mxCell>')
        for i, (s, t) in enumerate(edges):
            cells.append(f'<mxCell id="e{i}" edge="1" parent="1" source="{s}" '
                         'target="{}"><mxGeometry relative="1" as="geometry"/></mxCell>'.format(t))
        return ('<mxfile><diagram name="P"><mxGraphModel><root>'
                + "".join(cells) + "</root></mxGraphModel></diagram></mxfile>")

    @staticmethod
    def _status(style):
        for tag, hue in (("added", "82b366"), ("removed", "b85450"),
                         ("changed", "d79b00"), ("same", "999999")):
            if hue in style:
                return tag
        return "?"

    def test_drawiodiff_by_id(self):
        # api label moves (changed), cache removed, worker added, db unchanged.
        old = self._drawio([("api", "api v1"), ("db", "db"), ("cache", "cache")],
                           [("api", "db"), ("api", "cache")])
        new = self._drawio([("api", "api v2"), ("db", "db"), ("worker", "worker")],
                           [("api", "db"), ("api", "worker")])
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "old.drawio"), old)
            self._write(os.path.join(d, "new.drawio"), new)
            graph = json.loads(run("drawiodiff.py", os.path.join(d, "old.drawio"),
                                   os.path.join(d, "new.drawio")).stdout)
            status = {n["id"]: self._status(n["style"]) for n in graph["nodes"]}
            self.assertEqual(status, {"api": "changed", "db": "same",
                                      "cache": "removed", "worker": "added"})
            estatus = {(e["source"], e["target"]): self._status(e["style"])
                       for e in graph["edges"]}
            self.assertEqual(estatus, {("api", "db"): "same",
                                       ("api", "cache"): "removed",
                                       ("api", "worker"): "added"})

    def test_drawiodiff_by_label(self):
        # ids differ across versions; --by-label matches on the visible text.
        old = self._drawio([("x1", "API"), ("x2", "DB")], [])
        new = self._drawio([("y9", "API"), ("y8", "Cache")], [])
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "old.drawio"), old)
            self._write(os.path.join(d, "new.drawio"), new)
            base = run("drawiodiff.py", os.path.join(d, "old.drawio"),
                       os.path.join(d, "new.drawio"))
            # Without --by-label the random ids make everything add/remove.
            self.assertNotIn("same", {self._status(n["style"])
                                      for n in json.loads(base.stdout)["nodes"]})
            graph = json.loads(run("drawiodiff.py", os.path.join(d, "old.drawio"),
                                   os.path.join(d, "new.drawio"), "--by-label").stdout)
            status = {n["id"]: self._status(n["style"]) for n in graph["nodes"]}
            self.assertEqual(status, {"API": "same", "DB": "removed", "Cache": "added"})

    def test_timelapse_sample_indices(self):
        tl = load("timelapse")
        self.assertEqual(tl.sample_indices(3, 10), [0, 1, 2])        # fewer than n -> all
        self.assertEqual(tl.sample_indices(0, 5), [])
        picked = tl.sample_indices(20, 5)
        self.assertEqual(picked[0], 0)                                # always first
        self.assertEqual(picked[-1], 19)                             # always last
        self.assertEqual(len(picked), 5)
        self.assertEqual(picked, sorted(set(picked)))                # unique + ordered

    def test_timelapse_build_html_self_contained(self):
        tl = load("timelapse")
        frames = [(b"\x89PNG-old", "abc1234567", "2026-01-01", "init", 2, 1),
                  (b"\x89PNG-new", "def7654321", "2026-02-02", "grow", 9, 12)]
        html = tl.build_html(frames, "Evo")
        self.assertEqual(html.count("data:image/png;base64,"), 2)
        self.assertNotIn("http://", html)
        self.assertNotIn("https://", html)                           # no CDN/external refs
        for ctl in ('id="play"', 'id="scrub"', 'id="bar"'):
            self.assertIn(ctl, html)
        # the embedded frame payload carries per-commit counts
        self.assertIn('"n": 9', html)
        self.assertIn('"e": 12', html)

    # A container ("Tier") holding A + B, an ungrouped cylinder DB, one labeled edge.
    EXPLAIN_PAGE = (
        '<mxCell id="0"/><mxCell id="1" parent="0"/>'
        '<mxCell id="grp" value="Tier" vertex="1" parent="1" style="group">'
        '<mxGeometry x="0" y="0" width="240" height="120" as="geometry"/></mxCell>'
        '<mxCell id="a" value="A" vertex="1" parent="grp" style="rounded=1;">'
        '<mxGeometry x="10" y="30" width="80" height="40" as="geometry"/></mxCell>'
        '<mxCell id="b" value="B" vertex="1" parent="grp" style="rounded=1;">'
        '<mxGeometry x="120" y="30" width="80" height="40" as="geometry"/></mxCell>'
        '<mxCell id="c" value="DB" vertex="1" parent="1" style="shape=cylinder3;">'
        '<mxGeometry x="400" y="0" width="80" height="80" as="geometry"/></mxCell>'
        '<mxCell id="e1" value="reads" edge="1" parent="1" source="a" target="c">'
        '<mxGeometry relative="1" as="geometry"/></mxCell>')

    def test_explain_groups_shapes_relations(self):
        doc = ('<mxfile><diagram name="P1"><mxGraphModel><root>'
               + self.EXPLAIN_PAGE + "</root></mxGraphModel></diagram></mxfile>")
        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "arch.drawio")
            self._write(path, doc)
            md = run("explain.py", path).stdout
            self.assertIn("### Components (3)", md)       # A, B, DB (grp is a container)
            self.assertIn("- **Tier**", md)              # container -> grouping heading
            self.assertIn("- A", md)
            self.assertIn("DB _data store_", md)         # cylinder -> data store
            self.assertIn("A —reads→ DB", md)            # labeled relation
            self.assertNotIn("- Tier _", md)             # container itself isn't a component

    def test_explain_multipage(self):
        doc = ('<mxfile>'
               '<diagram name="P1"><mxGraphModel><root>' + self.EXPLAIN_PAGE
               + '</root></mxGraphModel></diagram>'
               '<diagram name="P2"><mxGraphModel><root>'
               '<mxCell id="0"/><mxCell id="1" parent="0"/>'
               '<mxCell id="x" value="Solo" vertex="1" parent="1" style="rounded=1;">'
               '<mxGeometry x="0" y="0" width="80" height="40" as="geometry"/></mxCell>'
               '</root></mxGraphModel></diagram></mxfile>')
        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "multi.drawio")
            self._write(path, doc)
            md = run("explain.py", path).stdout
            self.assertIn("## Page 1: P1", md)
            self.assertIn("## Page 2: P2", md)
            self.assertIn("- Solo", md)

    def test_drawio2pptx_page_names(self):
        tp = load("drawio2pptx")
        doc = ('<mxfile>'
               '<diagram name="Context"><mxGraphModel><root></root></mxGraphModel></diagram>'
               '<diagram name="Containers"><mxGraphModel><root></root></mxGraphModel></diagram>'
               '</mxfile>')
        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "c4.drawio")
            self._write(path, doc)
            self.assertEqual(tp.page_names(path), ["Context", "Containers"])

    def test_drawio2pptx_png_size(self):
        tp = load("drawio2pptx")
        import struct
        blob = b"\x00" * 16 + struct.pack(">II", 1280, 720)      # IHDR w/h at [16:24]
        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "x.png")
            with open(path, "wb") as f:
                f.write(blob)
            self.assertEqual(tp.png_size(path), (1280, 720))

    def test_drawio2pptx_deck(self):
        import shutil
        try:
            import pptx  # noqa: F401
        except ImportError:
            self.skipTest("python-pptx not installed")
        if not (shutil.which("drawio") and shutil.which("dot")):
            self.skipTest("draw.io CLI / Graphviz not installed")
        graph = {"direction": "TB",
                 "nodes": [{"id": "a", "label": "A"}, {"id": "b", "label": "B"}],
                 "edges": [{"source": "a", "target": "b"}]}
        with tempfile.TemporaryDirectory() as d:
            gj = os.path.join(d, "g.json")
            self._write(gj, json.dumps(graph))
            dr = os.path.join(d, "g.drawio")
            self.assertEqual(run("autolayout.py", gj, "-o", dr).returncode, 0)
            out = os.path.join(d, "g.pptx")
            r = run("drawio2pptx.py", dr, "-o", out)
            self.assertEqual(r.returncode, 0, r.stderr)
            from pptx import Presentation
            prs = Presentation(out)
            self.assertEqual(len(prs.slides), 1)

    # An edge line (pointer-events=stroke), an arrowhead + a shape (both =all).
    FLOW_SVG = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="100" height="50">'
        '<rect x="0" y="0" width="20" height="20" fill="#fff" pointer-events="all"/>'
        '<path d="M 20 10 L 60 10" fill="none" stroke="#000" pointer-events="stroke"/>'
        '<path d="M 60 10 L 55 7 L 55 13 Z" fill="#000000" pointer-events="all"/>'
        "</svg>")

    def test_svgflow_animates_only_edges(self):
        sf = load("svgflow")
        out, n = sf.animate(self.FLOW_SVG, 1.2, "6 4", reverse=False)
        self.assertEqual(n, 1)                                    # one edge line
        self.assertIn("@keyframes dio-flow", out)                # keyframes injected
        self.assertIn("stroke-dashoffset:-10", out)              # -(6+4) forward travel
        self.assertEqual(out.count('class="dio-flow"'), 1)
        # class is injected right after "<path " — only on the stroke edge
        self.assertIn('<path class="dio-flow" d="M 20 10', out)
        self.assertNotIn('<path class="dio-flow" d="M 60 10', out)  # arrowhead untouched
        self.assertNotIn('class="dio-flow"', out.split("<rect", 1)[1].split("/>", 1)[0])

    def test_svgflow_reverse_flips_offset(self):
        sf = load("svgflow")
        out, _ = sf.animate(self.FLOW_SVG, 2, "8 2", reverse=True)
        self.assertIn("stroke-dashoffset:10", out)               # +(8+2), toward source

    def test_drawio2mermaid_flowchart(self):
        # Reuses EXPLAIN_PAGE: container "Tier" with A/B, cylinder DB, edge a->c "reads".
        doc = ('<mxfile><diagram name="P1"><mxGraphModel><root>'
               + self.EXPLAIN_PAGE + "</root></mxGraphModel></diagram></mxfile>")
        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "arch.drawio")
            self._write(path, doc)
            mmd = run("drawio2mermaid.py", path, "--direction", "TD").stdout
            self.assertTrue(mmd.startswith("flowchart TD"))
            self.assertIn("subgraph", mmd)
            self.assertIn('["Tier"]', mmd)                       # container -> subgraph title
            self.assertIn('["A"]', mmd)
            self.assertIn('[("DB")]', mmd)                       # cylinder -> database shape
            self.assertIn('-->|"reads"|', mmd)                   # labeled edge

    def test_drawio2mermaid_multipage_fenced(self):
        doc = ('<mxfile>'
               '<diagram name="P1"><mxGraphModel><root>'
               '<mxCell id="0"/><mxCell id="1" parent="0"/>'
               '<mxCell id="x" value="Solo" vertex="1" parent="1" style="rounded=1;">'
               '<mxGeometry x="0" y="0" width="80" height="40" as="geometry"/></mxCell>'
               '</root></mxGraphModel></diagram>'
               '<diagram name="P2"><mxGraphModel><root>'
               '<mxCell id="0"/><mxCell id="1" parent="0"/>'
               '<mxCell id="y" value="Two" vertex="1" parent="1" style="rounded=1;">'
               '<mxGeometry x="0" y="0" width="80" height="40" as="geometry"/></mxCell>'
               '</root></mxGraphModel></diagram></mxfile>')
        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "multi.drawio")
            self._write(path, doc)
            mmd = run("drawio2mermaid.py", path, "--fenced").stdout
            self.assertEqual(mmd.count("```mermaid"), 2)          # one fence per page
            self.assertIn("%% Page 2: P2", mmd)
            self.assertIn('["Solo"]', mmd)

    SQL = ("CREATE TABLE users (id INT PRIMARY KEY, email VARCHAR(255));\n"
           "CREATE TABLE orders (\n"
           "  id INT,\n"
           "  user_id INT NOT NULL REFERENCES users(id),\n"
           "  PRIMARY KEY (id),\n"
           "  CONSTRAINT fk FOREIGN KEY (user_id) REFERENCES users (id)\n"
           ");\n")

    def test_sqlerd_tables_and_fks(self):
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "schema.sql"), self.SQL)
            graph = json.loads(run("sqlerd.py", d).stdout)
            ids = {n["id"] for n in graph["nodes"]}
            self.assertEqual(ids, {"users", "orders"})
            orders = next(n for n in graph["nodes"] if n["id"] == "orders")
            self.assertIn("PK id: INT", orders["label"])
            self.assertIn("FK user_id: INT", orders["label"])
            # inline REFERENCES + the equivalent named constraint -> 2 edges
            for e in graph["edges"]:
                self.assertEqual((e["source"], e["target"]), ("orders", "users"))
                self.assertIn("ERmany", e["style"])   # crow's foot at the FK side

    def test_seqlayout_geometry(self):
        spec = {"participants": [{"id": "a", "label": "A", "actor": True},
                                 {"id": "b", "label": "B"}],
                "messages": [{"from": "a", "to": "b", "label": "call"},
                             {"from": "b", "to": "b", "label": "work"},
                             {"from": "b", "to": "a", "label": "done", "return": True}]}
        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "seq.json")
            self._write(path, json.dumps(spec))
            out = os.path.join(d, "seq.drawio")
            r = run("seqlayout.py", path, "-o", out)
            self.assertEqual(r.returncode, 0, r.stderr)
            xml = open(out, encoding="utf-8").read()
            self.assertIn("umlActor", xml)                       # actor header
            self.assertIn('parent="b"', xml)                     # activation bar on b
            self.assertIn("endArrow=block", xml)                 # sync arrow
            self.assertIn("strokeColor=#999999", xml)            # return arrow
            # the generated file passes the structural linter
            v = run("validate.py", out)
            self.assertEqual(v.returncode, 0, v.stdout)

    def test_c4_multipage_drilldown(self):
        import shutil
        if not shutil.which("dot"):
            self.skipTest("Graphviz dot not installed")
        spec = {"levels": [
            {"name": "Context",
             "elements": [{"id": "u", "type": "person", "label": "User"},
                          {"id": "sys", "type": "system", "label": "System",
                           "children": "Containers"}],
             "relations": [{"from": "u", "to": "sys", "label": "Uses"}]},
            {"name": "Containers",
             "elements": [{"id": "api", "type": "container", "label": "API",
                           "tech": "Go"},
                          {"id": "db", "type": "database", "label": "DB"}],
             "relations": [{"from": "api", "to": "db"}]}]}
        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "c4.json")
            self._write(path, json.dumps(spec))
            out = os.path.join(d, "c4.drawio")
            r = run("c4.py", path, "-o", out)
            self.assertEqual(r.returncode, 0, r.stderr)
            xml = open(out, encoding="utf-8").read()
            self.assertEqual(xml.count("<diagram"), 2)
            self.assertIn('link="data:page/id,containers"', xml)   # drill-down
            self.assertIn("mxgraph.c4.person2", xml)
            self.assertIn("[Container: Go]", xml)
            v = run("validate.py", out)
            self.assertEqual(v.returncode, 0, v.stdout)

    def test_rustimports_edge(self):
        # Node ids are module paths (no crate-root file, so just a, b).
        with tempfile.TemporaryDirectory() as d:
            self._write(os.path.join(d, "Cargo.toml"), '[package]\nname = "m"\n')
            os.mkdir(os.path.join(d, "src"))
            self._write(os.path.join(d, "src", "a.rs"), "use crate::b;\n")
            self._write(os.path.join(d, "src", "b.rs"), "pub fn f() {}\n")
            r = run("rustimports.py", d)
            graph = json.loads(r.stdout)
            ids = {n["id"] for n in graph["nodes"]}
            self.assertEqual(len(graph["nodes"]), 2)
            self.assertEqual(ids, {"a", "b"})
            self.assertEqual(len(graph["edges"]), 1)
            self.assertEqual(graph["edges"][0], {"source": "a", "target": "b"})

    def test_openapiimports_operations_and_schemas(self):
        spec = {
            "openapi": "3.0.0",
            "paths": {"/pets": {
                "get": {"tags": ["pets"], "summary": "List",
                        "responses": {"200": {"content": {"application/json":
                            {"schema": {"$ref": "#/components/schemas/Pet"}}}}}},
                "post": {"tags": ["pets"],
                         "requestBody": {"content": {"application/json":
                            {"schema": {"$ref": "#/components/schemas/Pet"}}}},
                         "responses": {"201": {"description": "ok"}}},
            }},
            "components": {"schemas": {
                "Pet": {"properties": {"owner": {"$ref": "#/components/schemas/Owner"}}},
                "Owner": {"properties": {"id": {"type": "integer"}}},
            }},
        }
        with tempfile.TemporaryDirectory() as d:
            p = os.path.join(d, "api.json")
            self._write(p, json.dumps(spec))
            graph = json.loads(run("openapiimports.py", p, "--group").stdout)
            ops = [n for n in graph["nodes"] if n["id"].startswith("op")]
            schemas = [n for n in graph["nodes"] if n["id"].startswith("S:")]
            self.assertEqual(len(ops), 2)
            self.assertEqual({n["id"] for n in schemas}, {"S:Pet", "S:Owner"})
            get = next(n for n in ops if n["label"].startswith("GET"))
            post = next(n for n in ops if n["label"].startswith("POST"))
            self.assertIn("fillColor=#dae8fc", get["style"])       # GET blue
            self.assertIn("fillColor=#d5e8d4", post["style"])      # POST green
            pairs = {(e["source"], e["target"]) for e in graph["edges"]}
            self.assertIn((get["id"], "S:Pet"), pairs)             # operation -> schema
            self.assertIn(("S:Pet", "S:Owner"), pairs)            # schema nesting
            self.assertEqual({n.get("group") for n in ops}, {"pets"})
            self.assertEqual(next(iter(schemas)).get("group"), "schemas")

    def test_openapiimports_no_schemas(self):
        spec = {"openapi": "3.0.0",
                "paths": {"/x": {"get": {"responses": {"200": {"description": "ok"}}}}},
                "components": {"schemas": {"Y": {"properties": {"z": {"type": "string"}}}}}}
        with tempfile.TemporaryDirectory() as d:
            p = os.path.join(d, "api.json")
            self._write(p, json.dumps(spec))
            graph = json.loads(run("openapiimports.py", p, "--no-schemas").stdout)
            self.assertEqual(len(graph["nodes"]), 1)
            self.assertTrue(graph["nodes"][0]["id"].startswith("op"))
            self.assertEqual(graph["edges"], [])


class TestDrawioHtml(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.m = load("drawiohtml")

    TWO_PAGE = (
        '<mxfile>'
        '<diagram id="pg1" name="Context"><mxGraphModel><root>'
        '<mxCell id="0"/><mxCell id="1" parent="0"/>'
        '<UserObject label="System" link="data:page/id,pg2" id="sys">'
        '<mxCell style="rounded=1;" vertex="1" parent="1">'
        '<mxGeometry x="40" y="40" width="120" height="60" as="geometry"/></mxCell>'
        '</UserObject></root></mxGraphModel></diagram>'
        '<diagram id="pg2" name="Container"><mxGraphModel><root>'
        '<mxCell id="0"/><mxCell id="1" parent="0"/>'
        '<mxCell id="app" value="App" style="rounded=1;" vertex="1" parent="1">'
        '<mxGeometry x="40" y="40" width="120" height="60" as="geometry"/></mxCell>'
        '</root></mxGraphModel></diagram></mxfile>')

    def test_rewrite_page_links(self):
        import xml.etree.ElementTree as ET
        import io
        tree = ET.parse(io.StringIO(self.TWO_PAGE))
        self.assertEqual(self.m.rewrite_page_links(tree), 1)
        links = [el.get("link") for el in tree.getroot().iter() if el.get("link")]
        self.assertEqual(links, ["#page-pg2"])

    def test_strip_prolog(self):
        svg = '<?xml version="1.0"?>\n<!DOCTYPE svg>\n<svg><g/></svg>'
        self.assertEqual(self.m.strip_prolog(svg), "<svg><g/></svg>")

    def test_build_html_self_contained(self):
        html = self.m.build_html("t", [("pg1", "Context"), ("pg2", "Container")],
                                 ['<svg width="10" height="10"/>',
                                  '<svg width="10" height="10"/>'])
        self.assertEqual(html.count('class="page"'), 2)
        self.assertIn('data-pgid="pg1"', html)
        self.assertIn('"name": "Container"', html)
        self.assertIn("#page-", html)                          # link interception
        for banned in ("<script src=", "<link ", "@import", "https://cdn"):
            self.assertNotIn(banned, html)                     # no external requests

    def test_viewer_end_to_end(self):
        import shutil
        if not shutil.which("drawio"):
            self.skipTest("draw.io CLI not installed")
        with tempfile.TemporaryDirectory() as d:
            src = os.path.join(d, "two.drawio")
            with open(src, "w", encoding="utf-8") as f:
                f.write(self.TWO_PAGE)
            out = os.path.join(d, "two.html")
            r = run("drawiohtml.py", src, "-o", out)
            self.assertEqual(r.returncode, 0, r.stderr)
            self.assertIn("2 pages", r.stderr)
            self.assertIn("1 drill-down link", r.stderr)
            with open(out, encoding="utf-8") as f:
                html = f.read()
            self.assertEqual(html.count("<svg"), 2)            # both pages inlined
            self.assertIn("#page-pg2", html)                   # drill-down survived
            self.assertIn("data-cell-id", html)                # searchable cells


class TestHeatmap(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.m = load("heatmap")

    @staticmethod
    def _write(path, text):
        with open(path, "w", encoding="utf-8") as f:
            f.write(text)

    def test_ramp_hits_anchors(self):
        a = self.m.PALETTES["heat"]
        self.assertEqual(self.m.ramp(a, 0.0), "#57bb8a")          # low
        self.assertEqual(self.m.ramp(a, 0.5), "#ffd666")          # mid
        self.assertEqual(self.m.ramp(a, 1.0), "#e67c73")          # high

    def test_set_style_replaces_colours_once(self):
        s = self.m.set_style("rounded=1;fillColor=#111111;strokeColor=#222222;",
                             "#abcdef", "#123456")
        self.assertIn("fillColor=#abcdef", s)
        self.assertIn("strokeColor=#123456", s)
        self.assertNotIn("#111111", s)
        self.assertEqual(s.count("fillColor="), 1)

    def test_load_metrics_csv_and_json(self):
        with tempfile.TemporaryDirectory() as d:
            csvp = os.path.join(d, "m.csv")
            self._write(csvp, "service,ms\nweb,10\ndb,90\n")     # header auto-skipped
            self.assertEqual(self.m.load_metrics(csvp), {"web": 10.0, "db": 90.0})
            jp = os.path.join(d, "m.json")
            self._write(jp, '{"web":10,"db":90}')
            self.assertEqual(self.m.load_metrics(jp), {"web": 10.0, "db": 90.0})

    def test_recolors_matched_nodes(self):
        drawio = (
            '<mxfile><diagram id="p" name="P"><mxGraphModel><root>'
            '<mxCell id="0"/><mxCell id="1" parent="0"/>'
            '<mxCell id="web" value="Web" style="rounded=1;fillColor=#dddddd;'
            'strokeColor=#999999;" vertex="1" parent="1">'
            '<mxGeometry x="200" y="40" width="80" height="40" as="geometry"/></mxCell>'
            '<mxCell id="db" value="Database" style="rounded=1;fillColor=#dddddd;'
            'strokeColor=#999999;" vertex="1" parent="1">'
            '<mxGeometry x="360" y="40" width="80" height="40" as="geometry"/></mxCell>'
            '</root></mxGraphModel></diagram></mxfile>')
        with tempfile.TemporaryDirectory() as d:
            dp, mp, out = (os.path.join(d, n) for n in ("a.drawio", "m.csv", "out.drawio"))
            self._write(dp, drawio)
            self._write(mp, "web,10\ndb,90\n")
            r = run("heatmap.py", dp, "-m", mp, "-o", out)
            self.assertEqual(r.returncode, 0, r.stderr)
            self.assertIn("2/2 metrics matched", r.stderr)
            import xml.etree.ElementTree as ET
            cells = {c.get("id"): c.get("style") for c in ET.parse(out).getroot().iter("mxCell")}
            self.assertIn("fillColor=#57bb8a", cells["web"])       # min -> green
            self.assertIn("fillColor=#e67c73", cells["db"])        # max -> red
            self.assertIn("hm-title", cells)                       # legend injected


if __name__ == "__main__":
    unittest.main()
