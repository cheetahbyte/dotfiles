#!/usr/bin/env python3
"""Turn a .drawio into a PowerPoint deck — one slide per page.

Exports each page of a (multi-page) .drawio to a PNG via the draw.io CLI and
lays them into a 16:9 .pptx, one page per slide, scaled to fit with the page
name as the slide title. A C4 model (Context / Container / Component) becomes a
ready-to-present deck; any multi-page diagram becomes a slide sequence.

  python3 drawio2pptx.py c4.drawio -o c4.pptx
  python3 drawio2pptx.py architecture.drawio          # -> architecture.pptx

Needs the draw.io CLI (for the PNG export) and the `python-pptx` package
(`pip install python-pptx`) for writing the deck. Slides are 13.333in × 7.5in
(16:9); each image is centred and scaled to fit inside a small margin.

Usage: python3 drawio2pptx.py <file.drawio> [-o out.pptx] [--scale N]
"""
import argparse
import os
import struct
import subprocess
import sys
import tempfile
import xml.etree.ElementTree as ET


def page_names(path):
    """Names of the <diagram> pages, in order (None where unnamed)."""
    try:
        root = ET.parse(path).getroot()
    except (ET.ParseError, OSError) as exc:
        sys.exit(f"error: cannot parse {path}: {exc}")
    diagrams = root.findall("diagram")
    return [d.get("name") for d in diagrams] if diagrams else [None]


def png_size(path):
    """(width, height) in pixels from a PNG's IHDR header."""
    with open(path, "rb") as f:
        head = f.read(24)
    return struct.unpack(">II", head[16:24])


def export_page(drawio_file, index, out_png, scale):
    """Export one page (1-based index) to PNG via the draw.io CLI."""
    r = subprocess.run(["drawio", "-x", "-f", "png", "--page-index", str(index),
                        "-s", str(scale), "-o", out_png, drawio_file],
                       capture_output=True)
    return r.returncode == 0 and os.path.exists(out_png)


def main():
    ap = argparse.ArgumentParser(description="Export a .drawio to a PowerPoint deck (one slide per page).")
    ap.add_argument("file")
    ap.add_argument("-o", "--output", help="output .pptx (default: alongside input)")
    ap.add_argument("--scale", type=float, default=2.0, help="PNG export scale (default 2)")
    args = ap.parse_args()

    if not os.path.isfile(args.file):
        sys.exit(f"error: {args.file} not found")
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
    except ImportError:
        sys.exit("error: python-pptx is required (pip install python-pptx)")

    names = page_names(args.file)
    out = args.output or os.path.splitext(args.file)[0] + ".pptx"

    prs = Presentation()
    prs.slide_width = Emu(12192000)                    # 13.333in — 16:9
    prs.slide_height = Emu(6858000)                    # 7.5in
    blank = prs.slide_layouts[6]                        # the built-in "Blank" layout
    sw, sh = int(prs.slide_width), int(prs.slide_height)
    margin = Emu(457200)                               # 0.5in
    title_h = Emu(500000)

    made = 0
    with tempfile.TemporaryDirectory() as tmp:
        for i, name in enumerate(names, 1):            # draw.io --page-index is 1-based
            png = os.path.join(tmp, f"page{i}.png")
            if not export_page(args.file, i, png, args.scale):
                sys.stderr.write(f"warning: page {i} export failed — skipped\n")
                continue
            slide = prs.slides.add_slide(blank)
            top_pad = margin
            if name:
                box = slide.shapes.add_textbox(margin, Emu(180000),
                                               Emu(sw - 2 * int(margin)), title_h)
                tf = box.text_frame
                tf.text = name
                tf.paragraphs[0].runs[0].font.size = Pt(20)
                tf.paragraphs[0].runs[0].font.bold = True
                top_pad = Emu(180000) + title_h

            cw, ch = sw - 2 * int(margin), sh - int(top_pad) - int(margin)
            pw, ph = png_size(png)
            scale = min(cw / pw, ch / ph)              # fit, preserve aspect
            iw, ih = int(pw * scale), int(ph * scale)
            left = Emu(int((sw - iw) / 2))
            top = Emu(int(top_pad) + int((ch - ih) / 2))
            slide.shapes.add_picture(png, left, top, width=Emu(iw), height=Emu(ih))
            made += 1

    if not made:
        sys.exit("error: no pages exported (is the draw.io CLI installed?)")
    prs.save(out)
    sys.stderr.write(f"wrote {out} ({made} slide{'s' if made != 1 else ''})\n")


if __name__ == "__main__":
    main()
